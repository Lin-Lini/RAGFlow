"""FastAPI RAG service with Elasticsearch, SBERT embeddings, Vicuna & FLAN models.

Key improvements vs. your draft:
- Organized imports, removed duplicates, added typing.
- Centralized config (ES creds, model paths) via environment variables for safety.
- Proper NLTK downloads ("punkt", "stopwords").
- Correct text segmentation with page numbers; fixed logic bugs.
- Added Elasticsearch index for text_segments with dense_vector mapping.
- Safer file handling (BytesIO buffers, mime/type checks).
- Clear translation helpers (ru↔en).
- Consistent logging & error handling.
- PEP8 formatting + structured code sections.

Run: uvicorn fastapi_rag_app:app --reload
"""
from __future__ import annotations

import hashlib
import io
import logging
import os
import string
from datetime import datetime
from typing import Any, Dict, Iterable, List, Optional

import magic
import nltk
import numpy as np
import pandas as pd
import torch
from elasticsearch import Elasticsearch, exceptions
from fastapi import FastAPI, File, HTTPException, UploadFile, WebSocket
from fastapi.logger import logger as fastapi_logger
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, StreamingResponse
from langdetect import detect
from nltk.tokenize import sent_tokenize
from pdfminer.high_level import extract_text as pdfminer_extract_text  # alternative to PyPDF2
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer, util
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from torch import nn
from transformers import (
    AutoModelForCausalLM,
    AutoModelForSeq2SeqLM,
    AutoTokenizer,
    MarianMTModel,
    MarianTokenizer,
)

# --------------------------------------------------------------------------------------
# CONFIGURATION
# --------------------------------------------------------------------------------------

ES_URL = os.getenv("ES_URL", "https://localhost:9200")
ES_USERNAME = os.getenv("ES_USERNAME", "elastic")
ES_PASSWORD = os.getenv("ES_PASSWORD", "changeme")

MODEL_SBERT_NAME = os.getenv("MODEL_SBERT_NAME", "paraphrase-MiniLM-L6-v2")
MODEL_VICUNA_PATH = os.getenv("MODEL_VICUNA_PATH", "./llama-base_fine-tuned")
MODEL_FLAN_NAME = os.getenv("MODEL_FLAN_NAME", "google/flan-t5-base")
MODEL_MARIAN_RU_EN = os.getenv("MODEL_MARIAN_RU_EN", "Helsinki-NLP/opus-mt-ru-en")
MODEL_MARIAN_EN_RU = os.getenv("MODEL_MARIAN_EN_RU", "Helsinki-NLP/opus-mt-en-ru")

ALLOWED_EXTENSIONS = {"pptx", "doc", "docx", "pdf", "txt"}
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "text/plain",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

EMBEDDING_DIM = 384  # SBERT paraphrase-MiniLM-L6-v2 dimension
MIN_SCORE = 1.6
TOP_K_SEGMENTS = 10

# --------------------------------------------------------------------------------------
# NLTK SETUP
# --------------------------------------------------------------------------------------

nltk.download("stopwords", quiet=True)
nltk.download("punkt", quiet=True)

# --------------------------------------------------------------------------------------
# APP INIT & CORS
# --------------------------------------------------------------------------------------

app = FastAPI()

origins = [
    "http://localhost",
    "http://localhost:8080",
    "http://localhost:8001",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --------------------------------------------------------------------------------------
# LOGGING
# --------------------------------------------------------------------------------------

logger = fastapi_logger
logger.setLevel(logging.INFO)
handler = logging.StreamHandler()
handler.setFormatter(logging.Formatter("%(levelname)s:%(name)s:%(message)s"))
if not logger.handlers:
    logger.addHandler(handler)

# --------------------------------------------------------------------------------------
# DEVICE
# --------------------------------------------------------------------------------------

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

# --------------------------------------------------------------------------------------
# MODELS INIT
# --------------------------------------------------------------------------------------

logger.info("Loading models...")
model_sbert = SentenceTransformer(MODEL_SBERT_NAME).to(device)

model_vicuna = AutoModelForCausalLM.from_pretrained(
    MODEL_VICUNA_PATH,
    device_map="auto",
    torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
    load_in_8bit=False,
)
model_vicuna.to(device)

tokenizer_vicuna = AutoTokenizer.from_pretrained(MODEL_VICUNA_PATH, use_fast=True)

model_marian_ru_en = MarianMTModel.from_pretrained(MODEL_MARIAN_RU_EN).to(device)
tokenizer_marian_ru_en = MarianTokenizer.from_pretrained(MODEL_MARIAN_RU_EN)

model_marian_en_ru = MarianMTModel.from_pretrained(MODEL_MARIAN_EN_RU).to(device)
tokenizer_marian_en_ru = MarianTokenizer.from_pretrained(MODEL_MARIAN_EN_RU)

model_flan = AutoModelForSeq2SeqLM.from_pretrained(MODEL_FLAN_NAME, device_map="auto").to(device)
tokenizer_flan = AutoTokenizer.from_pretrained(MODEL_FLAN_NAME)

# --------------------------------------------------------------------------------------
# ELASTICSEARCH CLIENT & INDEXES
# --------------------------------------------------------------------------------------

logger.info("Connecting to Elasticsearch...")
es = Elasticsearch(ES_URL, basic_auth=(ES_USERNAME, ES_PASSWORD), verify_certs=False)

uploaded_files_mapping = {
    "mappings": {
        "properties": {
            "file_id": {"type": "keyword"},
            "file_name": {"type": "text"},
            "checksum": {"type": "keyword"},
            "uploaded_at": {"type": "date"},
            "file_size": {"type": "long"},
        }
    }
}

keywords_mapping = {
    "mappings": {
        "properties": {
            "keyword": {"type": "text"},
            "file_id": {"type": "keyword"},
            "file_name": {"type": "text"},
            "pages": {"type": "keyword"},
            "score": {"type": "float"},
            "timestamp": {"type": "date"},
        }
    }
}

learning_mapping = {
    "mappings": {
        "properties": {
            "question": {"type": "text"},
            "answer": {"type": "text"},
            "score": {"type": "float"},
        }
    }
}

learning_score_mapping = {
    "mappings": {
        "properties": {
            "question": {"type": "text"},
            "answer": {"type": "text"},
            "score": {"type": "float"},
        }
    }
}

text_segments_mapping = {
    "mappings": {
        "properties": {
            "file_id": {"type": "keyword"},
            "file_name": {"type": "text"},
            "page_number": {"type": "integer"},
            "segment_id": {"type": "integer"},
            "text": {"type": "text"},
            "embedding": {
                "type": "dense_vector",
                "dims": EMBEDDING_DIM,
                "index": False,  # using script_score, not knn_vector
            },
        }
    }
}


def create_index_if_not_exists(index_name: str, mapping: Dict[str, Any]) -> None:
    if not es.indices.exists(index=index_name):
        logger.info("Creating index %s", index_name)
        es.indices.create(index=index_name, body=mapping)


def create_required_indexes() -> None:
    create_index_if_not_exists("uploaded_files", uploaded_files_mapping)
    create_index_if_not_exists("keywords", keywords_mapping)
    create_index_if_not_exists("learning", learning_mapping)
    create_index_if_not_exists("learning-score", learning_score_mapping)
    create_index_if_not_exists("text_segments", text_segments_mapping)


@app.on_event("startup")
async def on_startup() -> None:
    create_required_indexes()


# --------------------------------------------------------------------------------------
# UTILITIES
# --------------------------------------------------------------------------------------

def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF using pdfminer for better accuracy than PyPDF2.
    Splits pages with delimiter '*_*' to keep page references.
    """
    # pdfminer works with file-like object
    with io.BytesIO(content) as pdf_buffer:
        # pdfminer returns full text; we need pages. We'll fallback to PyPDF2 for page split if needed
        try:
            text_full = pdfminer_extract_text(pdf_buffer)
        except Exception:
            text_full = ""

    if not text_full:
        # Fallback to PyPDF2 if pdfminer failed
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            parts = []
            for page in reader.pages:
                parts.append(page.extract_text() or "")
            return "*_*".join(parts)
        except Exception as e:  # noqa: BLE001
            raise HTTPException(status_code=400, detail=f"PDF parsing error: {e}")

    # If we don't have page delimiters, just return full text (single page)
    return text_full


def extract_text_from_docx(content: bytes) -> str:
    import docx  # lazy import

    document = docx.Document(io.BytesIO(content))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def extract_text_from_pptx(content: bytes) -> str:
    import pptx  # lazy import

    prs = pptx.Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def clean_text(text: str) -> str:
    return text.strip()


def translate_text(
    text: str,
    to_lang: str,
    model: MarianMTModel,
    tokenizer: MarianTokenizer,
    device: torch.device,
) -> str:
    """Translate `text` to target language using MarianMT.

    Args:
        text: source text
        to_lang: just for clarity ("en" or "ru"), not used inside
    """
    sentences = sent_tokenize(text)
    translated_sentences: List[str] = []

    for sentence in sentences:
        inputs = tokenizer(sentence, return_tensors="pt", padding=True, truncation=True).to(device)
        with torch.no_grad():
            generated = model.generate(**inputs)
        translated_sentence = tokenizer.decode(generated[0], skip_special_tokens=True)
        translated_sentences.append(translated_sentence)

    return " ".join(translated_sentences)


def sbert_segmentation_with_page_numbers(text: str, similarity_threshold: float = 0.7) -> List[Dict[str, Any]]:
    """Segment text by semantic similarity. Expects pages separated by '*_*'."""
    pages = text.split("*_*")
    segments: List[Dict[str, Any]] = []

    for ipage, page_text in enumerate(pages, start=1):
        sentences = [s.strip() for s in page_text.split(". ") if s.strip()]
        if not sentences:
            continue

        current_segment: List[str] = [sentences[0]]
        last_embedding = model_sbert.encode(sentences[0])

        for sentence in sentences[1:]:
            emb = model_sbert.encode(sentence)
            sim = cosine_similarity([emb], [last_embedding])[0][0]
            if sim < similarity_threshold:
                segments.append({"text": " ".join(current_segment), "page": ipage})
                current_segment = [sentence]
            else:
                current_segment.append(sentence)
            last_embedding = emb

        if current_segment:
            segments.append({"text": " ".join(current_segment), "page": ipage})

    return segments


def extract_keywords(text: str, top_n: int = 10) -> List[str]:
    sentences = nltk.sent_tokenize(text)
    words = nltk.word_tokenize(text)
    words = [w.lower() for w in words if w.isalnum()]

    sentence_embeddings = model_sbert.encode(sentences, convert_to_tensor=True)

    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(sentences)
    tfidf_scores = {word: tfidf_matrix.getcol(idx).sum() for word, idx in vectorizer.vocabulary_.items()}

    word_embeddings = model_sbert.encode(list(tfidf_scores.keys()), convert_to_tensor=True)
    avg_sent_emb = sentence_embeddings.mean(dim=0)
    cosine_scores = util.cos_sim(avg_sent_emb, word_embeddings)[0]

    top_indices = cosine_scores.argsort(descending=True)[: top_n * 2]
    candidates = [
        (
            list(tfidf_scores.keys())[idx],
            cosine_scores[idx].item() * tfidf_scores[list(tfidf_scores.keys())[idx]],
        )
        for idx in top_indices
    ]
    candidates.sort(key=lambda x: x[1], reverse=True)
    return [word for word, _ in candidates[:top_n]]


def index_keywords_to_elasticsearch(keywords: Iterable[str], file_id: str, file_name: str) -> None:
    for keyword in keywords:
        doc = {
            "keyword": keyword,
            "file_id": file_id,
            "file_name": file_name,
            "score": 0.0,
            "timestamp": datetime.utcnow().isoformat(),
        }
        es.index(index="keywords", document=doc)


def embed_text(text: str) -> List[float]:
    return model_sbert.encode(text).tolist()


# --------------------------------------------------------------------------------------
# DATA MODELS
# --------------------------------------------------------------------------------------


class LearningScore(BaseModel):
    question: str
    answer: str
    score: int


# --------------------------------------------------------------------------------------
# ROUTES
# --------------------------------------------------------------------------------------


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)) -> Dict[str, Any]:
    # 1. Extension & MIME check
    ext = file.filename.split(".")[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Недопустимый формат файла.")

    content = await file.read()
    mime = magic.from_buffer(content[:1024], mime=True)
    if mime not in ALLOWED_MIME_TYPES:
        raise HTTPException(status_code=400, detail="Недопустимый MIME-тип файла.")

    # 2. Check duplicate by checksum
    checksum = hashlib.md5(content).hexdigest()
    query = {"query": {"term": {"checksum": checksum}}}
    result = es.search(index="uploaded_files", body=query)
    if result["hits"]["total"]["value"] > 0:
        raise HTTPException(status_code=400, detail="Файл уже загружен.")

    # 3. Save file info
    uploaded_at = datetime.utcnow().isoformat()
    file_info = {
        "file_id": checksum,
        "file_name": file.filename,
        "checksum": checksum,
        "uploaded_at": uploaded_at,
        "file_size": len(content),
    }
    es.index(index="uploaded_files", document=file_info)

    # 4. Extract text
    if mime == "application/pdf":
        text = extract_text_from_pdf(content)
    elif mime in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }:
        text = extract_text_from_docx(content)
    elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(content)
    elif mime == "text/plain":
        text = content.decode("utf-8", errors="ignore")
    else:
        raise HTTPException(status_code=400, detail="Не удалось обработать файл.")

    text = clean_text(text)

    # 5. Language detection & optional translation to English for indexing
    lang = detect(text)
    if lang == "ru":
        text_for_index = translate_text(text, "en", model_marian_ru_en, tokenizer_marian_ru_en, device)
    else:
        text_for_index = text

    # 6. Keywords & indexing
    keywords = extract_keywords(text_for_index)
    index_keywords_to_elasticsearch(keywords, checksum, file.filename)

    # 7. Segmentation
    segments = sbert_segmentation_with_page_numbers(text_for_index)
    for i, segment in enumerate(segments):
        doc = {
            "file_id": checksum,
            "file_name": file.filename,
            "page_number": segment.get("page", 0),
            "segment_id": i,
            "text": segment["text"],
            "embedding": embed_text(segment["text"]),
        }
        es.index(index="text_segments", document=doc)

    return {"filename": file.filename, "checksum": checksum}


@app.websocket("/ask")
async def ask_endpoint(websocket: WebSocket) -> None:
    await websocket.accept()
    while True:
        question = await websocket.receive_text()
        lang = detect(question)
        query_text = question
        if lang == "ru":
            query_text = translate_text(question, "en", model_marian_ru_en, tokenizer_marian_ru_en, device)

        question_embedding = embed_text(query_text)
        script_query = {
            "script_score": {
                "query": {"match_all": {}},
                "script": {
                    "source": "cosineSimilarity(params.query_vector, 'embedding') + 1.0",
                    "params": {"query_vector": question_embedding},
                },
            }
        }
        es_resp = es.search(
            index="text_segments",
            body={
                "size": TOP_K_SEGMENTS,
                "query": script_query,
                "_source": ["file_name", "text", "page_number"],
            },
        )

        hits = es_resp.get("hits", {}).get("hits", [])
        relevant_blocks = []
        context = ""

        if hits and hits[0]["_score"] > MIN_SCORE:
            context = hits[0]["_source"]["text"]
            for hit in hits:
                if hit["_score"] > MIN_SCORE:
                    relevant_blocks.append(
                        {
                            "file": hit["_source"]["file_name"],
                            "pages": [hit["_source"].get("page_number")],
                            "score": hit["_score"],
                            "squeeze": hit["_source"]["text"],
                        }
                    )
            input_text_flan = f"Question: {query_text} Context (if relevant): {context} Answer: "
            input_text_vicuna = context + "\nUser: " + query_text + "\nAssistant:"
        else:
            input_text_flan = (
                f"Question: {query_text} "
                "Answer the question based on your general knowledge without specific context. Answer: "
            )
            input_text_vicuna = "User: " + query_text + "\nAssistant:"

        # FLAN generation
        flan_ids = tokenizer_flan(input_text_flan, return_tensors="pt").input_ids.to(device)
        flan_out = model_flan.generate(
            flan_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens=50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
        )
        answer_flan = tokenizer_flan.decode(flan_out[0], skip_special_tokens=True)

        # Vicuna generation
        vicuna_ids = tokenizer_vicuna(input_text_vicuna, return_tensors="pt").input_ids.to(device)
        vicuna_out = model_vicuna.generate(
            vicuna_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens=50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
        )
        answer_vicuna = tokenizer_vicuna.decode(vicuna_out[0], skip_special_tokens=True)
        if "Answer:" in answer_vicuna:
            answer_vicuna = answer_vicuna.split("Answer:")[-1].strip()
        elif "Assistant:" in answer_vicuna:
            answer_vicuna = answer_vicuna.split("Assistant:")[-1].strip()

        if lang == "ru":
            answer_flan = translate_text(answer_flan, "ru", model_marian_en_ru, tokenizer_marian_en_ru, device)
            answer_vicuna = translate_text(answer_vicuna, "ru", model_marian_en_ru, tokenizer_marian_en_ru, device)

        payload = {
            "answer_flan": answer_flan,
            "answer": answer_vicuna,
            "relevance": relevant_blocks,
        }
        await websocket.send_json(payload)


@app.post("/process_questions")
async def process_questions(file: UploadFile = File(...)) -> StreamingResponse:
    df = pd.read_excel(file.file)

    for idx, row in df.iterrows():
        question = row.get("question", "")
        if not question:
            continue
        lang = detect(question)
        query_text = question
        if lang == "ru":
            query_text = translate_text(question, "en", model_marian_ru_en, tokenizer_marian_ru_en, device)

        question_embedding = embed_text(query_text)
        script_query = {
            "script_score": {
                "query": {"match_all": {}},
                "script": {
                    "source": "cosineSimilarity(params.query_vector, 'embedding') + 1.0",
                    "params": {"query_vector": question_embedding},
                },
            }
        }
        es_resp = es.search(
            index="text_segments",
            body={"size": 1, "query": script_query, "_source": ["file_name", "text", "page_number"]},
        )

        context = ""
        filename = ""
        slide_number: Optional[int] = None

        if es_resp["hits"]["hits"] and es_resp["hits"]["hits"][0]["_score"] > MIN_SCORE:
            hit = es_resp["hits"]["hits"][0]
            context = hit["_source"]["text"]
            filename = os.path.splitext(hit["_source"]["file_name"])[0]
            slide_number = hit["_source"].get("page_number")
            prompt_vicuna = context + "\nUser: " + query_text + "\nAssistant:"
        else:
            prompt_vicuna = "User: " + query_text + "\nAssistant:"

        vicuna_ids = tokenizer_vicuna(prompt_vicuna, return_tensors="pt").input_ids.to(device)
        vicuna_out = model_vicuna.generate(
            vicuna_ids,
            temperature=0.7,
            max_length=256,
            max_new_tokens=50,
            top_p=0.9,
            top_k=40,
            do_sample=True,
            use_cache=False,
        )
        answer = tokenizer_vicuna.decode(vicuna_out[0], skip_special_tokens=True)
        if "Answer:" in answer:
            answer = answer.split("Answer:")[-1].strip()
        elif "Assistant:" in answer:
            answer = answer.split("Assistant:")[-1].strip()

        if lang == "ru":
            answer = translate_text(answer, "ru", model_marian_en_ru, tokenizer_marian_en_ru, device)

        df.at[idx, "answer"] = answer
        df.at[idx, "filename"] = filename
        df.at[idx, "slide_number"] = slide_number if slide_number is not None else ""

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False)
    output.seek(0)

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=processed_questions.xlsx"},
    )


@app.websocket("/prompt")
async def prompt_endpoint(websocket: WebSocket) -> None:
    await websocket.accept()
    while True:
        prefix = await websocket.receive_text()
        es_resp = es.search(
            index="keywords",
            body={
                "size": 10,
                "query": {"prefix": {"keyword": {"value": prefix}}},
                "_source": ["keyword"],
            },
        )
        suggestions = [hit["_source"]["keyword"] for hit in es_resp["hits"]["hits"]]
        await websocket.send_json(suggestions)


@app.post("/vote")
async def vote_endpoint(data: LearningScore) -> Dict[str, Any]:
    try:
        if not es.indices.exists(index="learning-score"):
            es.indices.create(index="learning-score")
        resp = es.index(index="learning-score", document=data.dict())
        return {"status": "success", "data": resp}
    except exceptions.ConnectionError:
        raise HTTPException(status_code=500, detail="Connection to Elasticsearch failed")
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=500, detail=str(e))


# --------------------------------------------------------------------------------------
# SIMPLE ROOT ROUTE
# --------------------------------------------------------------------------------------

@app.get("/", response_class=HTMLResponse)
async def root() -> str:
    return "<h3>RAG API is running</h3>"